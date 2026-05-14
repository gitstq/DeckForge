"""
PPTX Renderer - Render slide data to PPTX format
PPTX渲染器 - 将幻灯片数据渲染为PPTX格式
"""

from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class PPTXRenderer:
    """Render structured slide data to PPTX format."""
    
    # Standard slide dimensions (16:9)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)
    
    def __init__(self, theme: Dict[str, Any]):
        """Initialize renderer with theme.
        
        Args:
            theme: Theme configuration dictionary
        """
        self.theme = theme
        self.colors = theme.get('colors', {})
        self.fonts = theme.get('fonts', {})
        self.layouts = theme.get('layouts', {})
    
    def render(self, slides: List[Dict[str, Any]]) -> Presentation:
        """Render slides to a PPTX Presentation object.
        
        Args:
            slides: List of slide dictionaries
            
        Returns:
            pptx.Presentation object
        """
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT
        
        for slide_data in slides:
            self._render_slide(prs, slide_data)
        
        return prs
    
    def _render_slide(self, prs: Presentation, slide_data: Dict[str, Any]) -> None:
        """Render a single slide.
        
        Args:
            prs: Presentation object
            slide_data: Slide dictionary
        """
        slide_type = slide_data.get('type', 'content')
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Apply background
        self._apply_background(slide)
        
        # Render by type
        if slide_type == 'title':
            self._render_title_slide(slide, slide_data)
        elif slide_type == 'section':
            self._render_section_slide(slide, slide_data)
        elif slide_type == 'closing':
            self._render_closing_slide(slide, slide_data)
        elif slide_type == 'bullets':
            self._render_bullets_slide(slide, slide_data)
        elif slide_type == 'content':
            self._render_content_slide(slide, slide_data)
        elif slide_type == 'comparison':
            self._render_comparison_slide(slide, slide_data)
        elif slide_type == 'two_column':
            self._render_two_column_slide(slide, slide_data)
        elif slide_type == 'code':
            self._render_code_slide(slide, slide_data)
        elif slide_type == 'image':
            self._render_content_slide(slide, slide_data)
        elif slide_type == 'table':
            self._render_content_slide(slide, slide_data)
        else:
            self._render_content_slide(slide, slide_data)
        
        # Add speaker notes
        notes = slide_data.get('notes', '')
        if notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes
    
    def _apply_background(self, slide) -> None:
        """Apply theme background to slide.
        
        Args:
            slide: Slide object
        """
        bg_color = self.colors.get('background', '#FFFFFF')
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(bg_color.lstrip('#'))
    
    def _add_text_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_size: int = 18,
        font_color: str = '#333333',
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
        font_name: Optional[str] = None
    ):
        """Add a text box to a slide.
        
        Args:
            slide: Slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            text: Text content
            font_size: Font size in points
            font_color: Hex color string
            bold: Whether text is bold
            alignment: Text alignment
            font_name: Custom font name
            
        Returns:
            Text box shape
        """
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(font_color.lstrip('#'))
        p.font.bold = bold
        p.font.name = font_name or self.fonts.get('body', 'Arial')
        p.alignment = alignment
        
        return txBox
    
    def _add_shape_bg(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        color: str,
        alpha: Optional[int] = None
    ):
        """Add a colored rectangle shape as background element.
        
        Args:
            slide: Slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            color: Hex color string
            alpha: Opacity percentage (0-100)
        """
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color.lstrip('#'))
        shape.line.fill.background()
    
    def _render_title_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a title slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', 'Presentation')
        subtitle = data.get('subtitle', '')
        
        # Accent bar at top
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 0.15, accent_color)
        
        # Title
        title_color = self.colors.get('title', '#1a1a2e')
        title_font = self.fonts.get('title', 'Arial')
        
        self._add_text_box(
            slide, 1.0, 2.0, 11.333, 2.0,
            title,
            font_size=44,
            font_color=title_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_name=title_font
        )
        
        # Subtitle
        if subtitle:
            sub_color = self.colors.get('subtitle', '#666666')
            self._add_text_box(
                slide, 1.5, 4.2, 10.333, 1.5,
                subtitle,
                font_size=20,
                font_color=sub_color,
                alignment=PP_ALIGN.CENTER,
                font_name=self.fonts.get('body', 'Arial')
            )
        
        # Bottom accent line
        self._add_shape_bg(slide, 5.5, 6.5, 2.333, 0.05, accent_color)
        
        # Branding
        brand_color = self.colors.get('muted', '#999999')
        self._add_text_box(
            slide, 0.5, 6.8, 12.333, 0.5,
            'Created with DeckForge',
            font_size=10,
            font_color=brand_color,
            alignment=PP_ALIGN.CENTER
        )
    
    def _render_section_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a section divider slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', 'Section')
        
        # Full-width accent background
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 2.5, 13.333, 2.5, accent_color)
        
        # Section title
        self._add_text_box(
            slide, 1.0, 3.0, 11.333, 1.5,
            title,
            font_size=40,
            font_color='#FFFFFF',
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_name=self.fonts.get('title', 'Arial')
        )
    
    def _render_closing_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a closing/thank you slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', 'Thank You')
        subtitle = data.get('subtitle', '')
        
        # Accent bar
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 0.15, accent_color)
        
        # Thank you text
        self._add_text_box(
            slide, 1.0, 2.5, 11.333, 2.0,
            title,
            font_size=48,
            font_color=self.colors.get('title', '#1a1a2e'),
            bold=True,
            alignment=PP_ALIGN.CENTER,
            font_name=self.fonts.get('title', 'Arial')
        )
        
        if subtitle:
            self._add_text_box(
                slide, 1.5, 4.5, 10.333, 1.0,
                subtitle,
                font_size=22,
                font_color=self.colors.get('subtitle', '#666666'),
                alignment=PP_ALIGN.CENTER
            )
    
    def _render_bullets_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a bullet points slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', '')
        bullets = data.get('bullets', [])
        
        # Header bar
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 1.2, accent_color)
        
        # Title
        self._add_text_box(
            slide, 0.8, 0.15, 11.733, 0.9,
            title,
            font_size=28,
            font_color='#FFFFFF',
            bold=True,
            font_name=self.fonts.get('title', 'Arial')
        )
        
        # Bullets
        if bullets:
            y_pos = 1.6
            bullet_color = self.colors.get('text', '#333333')
            
            for bullet in bullets[:7]:  # Max 7 bullets per slide
                # Bullet marker
                self._add_shape_bg(
                    slide, 1.0, y_pos + 0.15, 0.15, 0.15,
                    accent_color
                )
                
                # Bullet text
                self._add_text_box(
                    slide, 1.4, y_pos, 10.933, 0.6,
                    bullet,
                    font_size=18,
                    font_color=bullet_color,
                    font_name=self.fonts.get('body', 'Arial')
                )
                y_pos += 0.7
    
    def _render_content_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a content slide with title and body text.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', '')
        content = data.get('content', '')
        
        # Header bar
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 1.2, accent_color)
        
        # Title
        self._add_text_box(
            slide, 0.8, 0.15, 11.733, 0.9,
            title,
            font_size=28,
            font_color='#FFFFFF',
            bold=True,
            font_name=self.fonts.get('title', 'Arial')
        )
        
        # Body content
        if content:
            # Truncate if too long
            display_content = content[:500] if len(content) > 500 else content
            
            txBox = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.4)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            
            lines = display_content.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = line.strip()
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor.from_string(
                    self.colors.get('text', '#333333').lstrip('#')
                )
                p.font.name = self.fonts.get('body', 'Arial')
                p.space_after = Pt(8)
    
    def _render_comparison_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a comparison/two-column slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', '')
        content = data.get('content', '')
        
        # Header
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 1.2, accent_color)
        
        self._add_text_box(
            slide, 0.8, 0.15, 11.733, 0.9,
            title,
            font_size=28,
            font_color='#FFFFFF',
            bold=True,
            font_name=self.fonts.get('title', 'Arial')
        )
        
        # Two columns
        lines = [l.strip() for l in content.split('\n') if l.strip()]
        mid = len(lines) // 2
        
        left_text = '\n'.join(lines[:mid]) if mid > 0 else content
        right_text = '\n'.join(lines[mid:]) if mid > 0 else ''
        
        # Left column background
        self._add_shape_bg(
            slide, 0.5, 1.5, 5.9, 5.5,
            self.colors.get('surface', '#F5F5F5')
        )
        
        # Right column background
        self._add_shape_bg(
            slide, 6.9, 1.5, 5.9, 5.5,
            self.colors.get('surface', '#F5F5F5')
        )
        
        # Left text
        self._add_text_box(
            slide, 0.8, 1.7, 5.3, 5.0,
            left_text[:300],
            font_size=15,
            font_color=self.colors.get('text', '#333333'),
            font_name=self.fonts.get('body', 'Arial')
        )
        
        # Right text
        if right_text:
            self._add_text_box(
                slide, 7.2, 1.7, 5.3, 5.0,
                right_text[:300],
                font_size=15,
                font_color=self.colors.get('text', '#333333'),
                font_name=self.fonts.get('body', 'Arial')
            )
    
    def _render_two_column_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a two-column layout slide."""
        self._render_comparison_slide(slide, data)
    
    def _render_code_slide(self, slide, data: Dict[str, Any]) -> None:
        """Render a code display slide.
        
        Args:
            slide: Slide object
            data: Slide data
        """
        title = data.get('title', '')
        content = data.get('content', '')
        
        # Header
        accent_color = self.colors.get('accent', '#4A90D9')
        self._add_shape_bg(slide, 0, 0, 13.333, 1.2, accent_color)
        
        self._add_text_box(
            slide, 0.8, 0.15, 11.733, 0.9,
            title,
            font_size=28,
            font_color='#FFFFFF',
            bold=True,
            font_name=self.fonts.get('title', 'Arial')
        )
        
        # Code block background
        self._add_shape_bg(
            slide, 0.8, 1.5, 11.733, 5.5,
            self.colors.get('code_bg', '#1e1e1e')
        )
        
        # Code text
        # Extract code from markdown code blocks
        code_match = content.find('```')
        if code_match != -1:
            start = content.find('\n', code_match) + 1
            end = content.find('```', start)
            if end != -1:
                code_text = content[start:end]
            else:
                code_text = content[code_match + 3:]
        else:
            code_text = content
        
        # Truncate long code
        code_text = code_text[:800] if len(code_text) > 800 else code_text
        
        txBox = slide.shapes.add_textbox(
            Inches(1.2), Inches(1.8), Inches(10.933), Inches(5.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(code_text.split('\n')[:20]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(13)
            p.font.color.rgb = RGBColor.from_string('#D4D4D4')
            p.font.name = self.fonts.get('code', 'Consolas')
            p.space_after = Pt(2)
